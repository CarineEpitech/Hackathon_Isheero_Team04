#!/usr/bin/env python3
# Génère TERROIR_pitch_deck_v5.pptx dans C:/Users/carin/OneDrive/Downloads/
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT  = r"C:\Users\carin\OneDrive\Downloads"
URL  = "terroir.up.railway.app/#/live"
TEAM = "BeninScope"

# ── Palette ───────────────────────────────────────────────────────────────────
BG   = RGBColor(0x10, 0x10, 0x10)   # noir chaud
PH   = RGBColor(0x1A, 0x1A, 0x1A)   # placeholder
PHB  = RGBColor(0x40, 0x40, 0x40)   # bordure placeholder
GRN  = RGBColor(0x22, 0xC5, 0x5E)   # vert
AMB  = RGBColor(0xEA, 0xB3, 0x08)   # ambre
RED  = RGBColor(0xEF, 0x44, 0x44)   # rouge
WHT  = RGBColor(0xF8, 0xF8, 0xF8)   # blanc
MUT  = RGBColor(0x72, 0x72, 0x72)   # gris muet

# ── Helpers de base ──────────────────────────────────────────────────────────

def new_prs():
    p = Presentation()
    p.slide_width  = In(13.33)
    p.slide_height = In(7.5)
    return p

def slide(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    f = s.background.fill; f.solid()
    f.fore_color.rgb = BG
    return s

def box(s, l, t, w, h, fill=None, lc=None, lw=1.0):
    sh = s.shapes.add_shape(1, In(l), In(t), In(w), In(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:    sh.fill.background()
    if lc:   sh.line.color.rgb = lc; sh.line.width = Pt(lw)
    else:    sh.line.fill.background()
    return sh

def tx(s, l, t, w, h, text, sz=16, bold=False, col=None,
       al=PP_ALIGN.LEFT, it=False):
    col = col or WHT
    b   = s.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf  = b.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = al
    r   = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = it;   r.font.color.rgb = col
    return b

def txm(s, l, t, w, h, lines):
    b   = s.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf  = b.text_frame; tf.word_wrap = True
    fst = True
    for item in lines:
        if fst: p = tf.paragraphs[0]; fst = False
        else:   p = tf.add_paragraph()
        if not item: continue
        if isinstance(item, str): item = {'text': item}
        p.alignment = item.get('al', PP_ALIGN.LEFT)
        sp = item.get('sp', 0)
        if sp:
            from pptx.util import Pt as _Pt
            p.space_before = _Pt(sp)
        r = p.add_run()
        r.text          = item.get('text', '')
        r.font.size     = Pt(item.get('sz', 16))
        r.font.bold     = item.get('bold', False)
        r.font.italic   = item.get('it', False)
        r.font.color.rgb = item.get('col', WHT)
    return b

def vbar(s, col=None, w=0.07):
    """Barre verticale gauche — ancre visuelle sobre."""
    box(s, 0, 0.55, w, 6.4, fill=col or GRN)

def sn(s, n, total=7):
    tx(s, 12.6, 7.18, 0.6, 0.22, f"{n}/{total}", sz=9, col=MUT, al=PP_ALIGN.RIGHT)

def capture_ph(s, l, t, w, h, title_label, detail=''):
    """Placeholder visuel pour capture d'écran."""
    box(s, l, t, w, h, fill=PH, lc=PHB, lw=0.75)
    # Texte centré dans la zone
    cy = t + (h / 2) - 0.55
    tx(s, l + 0.2, cy, w - 0.4, 0.45,
       "[ CAPTURE À INSÉRER ]",
       sz=11, col=PHB, al=PP_ALIGN.CENTER)
    tx(s, l + 0.2, cy + 0.5, w - 0.4, 0.45,
       title_label,
       sz=13, bold=True, col=MUT, al=PP_ALIGN.CENTER)
    if detail:
        tx(s, l + 0.2, cy + 1.0, w - 0.4, 0.5,
           detail,
           sz=10, col=MUT, al=PP_ALIGN.CENTER, it=True)

def qr_ph(s, l, t, w=1.5, h=1.5):
    """Placeholder QR code."""
    box(s, l, t, w, h, fill=WHT, lc=MUT, lw=0.5)
    tx(s, l + 0.05, t + 0.05, w - 0.1, h - 0.1,
       "QR CODE\nterroir.up.railway.app",
       sz=9, col=MUT, al=PP_ALIGN.CENTER)

# ── SLIDES ────────────────────────────────────────────────────────────────────

def s1_cover(p):
    s = slide(p)
    # Barre verte verticale gauche
    vbar(s)
    # Titre principal
    tx(s, 0.55, 1.4, 9.0, 1.8, "TERROIR", sz=66, bold=True, col=WHT)
    # Ligne verte sous titre
    box(s, 0.55, 3.35, 2.2, 0.06, fill=GRN)
    # Sous-titre
    tx(s, 0.55, 3.6, 9.5, 0.65,
       "Veille sécuritaire territoriale · Bénin · Temps réel",
       sz=18, col=MUT)
    # URL
    tx(s, 0.55, 4.5, 8.0, 0.45, URL, sz=13, col=GRN)
    # Séparation bas
    box(s, 0.55, 6.55, 12.2, 0.04, fill=MUT)
    # Équipe
    tx(s, 0.55, 6.7, 6.0, 0.45, TEAM + "  ·  iSHEERO × DataCamp 2026",
       sz=11, col=MUT)
    # Membres
    tx(s, 0.55, 7.1, 10.0, 0.3,
       "Yaoitcha Rosine  ·  Gbaguidi Adewale  ·  Houndofi Jacques  ·  Agboton Carine",
       sz=9, col=MUT)
    # Petits carrés décoratifs (carte symbolique, discrets)
    for x, y, c in [(11.8, 2.0, GRN), (12.3, 3.2, AMB), (11.5, 4.1, MUT)]:
        box(s, x, y, 0.12, 0.12, fill=c)

def s2_probleme(p):
    s = slide(p)
    vbar(s)
    sn(s, 2)
    # Phrase forte — grande
    tx(s, 0.55, 0.8, 11.5, 1.3,
       "Le nord du Bénin est sous tension.",
       sz=36, bold=True, col=WHT)
    tx(s, 0.55, 2.25, 9.0, 0.65,
       "Peu couvert par les médias. Mais réel.",
       sz=20, col=MUT)
    # Séparateur
    box(s, 0.55, 3.1, 1.5, 0.05, fill=GRN)
    # 2 chiffres — pas dans des boîtes, juste typographie
    txm(s, 0.55, 3.4, 6.5, 2.8, [
        {'text': "4 %",  'sz': 52, 'bold': True, 'col': AMB},
        {'text': "de la couverture médiatique internationale du pays",
         'sz': 14, 'col': MUT},
        {'text': " ", 'sz': 10, 'col': MUT},
        {'text': "4 ×", 'sz': 52, 'bold': True, 'col': RED},
        {'text': "plus négatif que le sud  (ton −4,10 vs −1,09)",
         'sz': 14, 'col': MUT},
    ])
    # Note à droite — ce que ça veut dire
    tx(s, 7.5, 3.4, 5.3, 3.0,
       "Quand la zone nord est couverte, "
       "c'est presque toujours pour un incident "
       "de sécurité interne. Les signaux existent. "
       "Ils ne sont pas agrégés.",
       sz=14, col=MUT, it=True)

def s3_outils(p):
    s = slide(p)
    vbar(s, col=RED)
    sn(s, 3)
    # Titre sobre
    tx(s, 0.55, 0.7, 11.0, 0.9,
       "Les outils existants arrivent trop tard.",
       sz=30, bold=True, col=WHT)
    box(s, 0.55, 1.7, 1.2, 0.05, fill=RED)
    # 3 points — texte brut, pas de boîtes
    items = [
        ("ACLED",                    "Enregistre les incidents après qu'ils se produisent."),
        ("Rapports institutionnels", "Hebdomadaires ou mensuels. L'urgence n'attend pas."),
        ("Dataminr / ICEWS",         "Conçus pour les gouvernements occidentaux. Inaccessibles."),
    ]
    for i, (label, desc) in enumerate(items):
        y = 2.15 + i * 1.5
        # Tiret vert
        box(s, 0.55, y + 0.2, 0.25, 0.05, fill=RED)
        tx(s, 0.9, y, 4.0, 0.5, label, sz=18, bold=True, col=WHT)
        tx(s, 0.9, y + 0.5, 11.0, 0.6, desc, sz=14, col=MUT)

def s4_carte(p):
    s = slide(p)
    sn(s, 4)
    # Titre court en haut à gauche
    tx(s, 0.5, 0.25, 7.0, 0.65, "La carte live.", sz=24, bold=True, col=WHT)
    tx(s, 0.5, 0.95, 9.0, 0.4,
       "GDELT · Signalements terrain · Temps réel",
       sz=13, col=MUT)
    # Capture principale — grande
    capture_ph(s, 0.4, 1.5, 12.5, 5.5,
               "Carte Live — terroir.up.railway.app/#/live",
               "Montrer : carte Bénin avec marqueurs événements + signalements terrain")
    # URL bas
    tx(s, 0.5, 7.1, 8.0, 0.3, URL, sz=10, col=GRN)

def s5_comment(p):
    s = slide(p)
    vbar(s, col=GRN)
    sn(s, 5)
    tx(s, 0.55, 0.7, 9.0, 0.8,
       "Deux sources. Un signal. Une décision.",
       sz=28, bold=True, col=WHT)
    box(s, 0.55, 1.6, 1.5, 0.05, fill=GRN)
    # Flow vertical — texte brut, sobre
    steps = [
        (GRN, "GDELT",              "Médias internationaux · toutes les 15 min"),
        (GRN, "Signalements terrain","Citoyens géolocalisés · temps réel"),
        (AMB, "Carte live",         "Tous les acteurs voient la même carte"),
        (AMB, "Score de tension",   "Par département · calculé quotidiennement"),
        (WHT, "Alerte",             "Zone · score · précédent · recommandation"),
        (WHT, "Décision",           "Humaine · tracée · documentée"),
    ]
    # Colonne gauche : labels + descriptions
    # Colonne droite : rien (respiration)
    y = 1.85
    for i, (col, label, desc) in enumerate(steps):
        # Flèche sauf premier
        if i > 0:
            tx(s, 1.5, y - 0.32, 0.4, 0.35, "↓", sz=14, col=MUT, al=PP_ALIGN.CENTER)
        box(s, 0.55, y, 0.08, 0.38, fill=col)
        tx(s, 0.75, y, 2.8, 0.4, label, sz=15, bold=True, col=col)
        tx(s, 3.7,  y, 8.5, 0.4, desc,  sz=13, col=MUT)
        y += 0.72

def s6_impact(p):
    s = slide(p)
    vbar(s, col=AMB)
    sn(s, 6)
    # Titre
    tx(s, 0.55, 0.7, 12.0, 0.7, "Pour qui. Pourquoi c'est différent.", sz=26, bold=True)
    box(s, 0.55, 1.5, 1.2, 0.05, fill=AMB)

    # Gauche : utilisateurs
    tx(s, 0.55, 1.75, 3.0, 0.45, "Qui l'utilise", sz=13, bold=True, col=GRN)
    users = ["ONG", "Mairie", "Journaliste", "Préfecture"]
    for i, u in enumerate(users):
        y = 2.35 + i * 0.9
        box(s, 0.55, y, 0.07, 0.42, fill=GRN)
        tx(s, 0.75, y, 3.5, 0.45, u, sz=18, bold=True, col=WHT)

    # Séparateur vertical
    box(s, 4.6, 1.75, 0.04, 4.8, fill=MUT)

    # Droite : différenciation
    tx(s, 4.9, 1.75, 8.0, 0.45, "Ce que les autres ne font pas", sz=13, bold=True, col=AMB)
    diffs = [
        ("Signal avant l'incident",
         "Pas après. TERROIR détecte les précurseurs médiatiques."),
        ("Terrain + médias combinés",
         "GDELT + signalements GPS. La couche citoyenne corrige le biais GDELT."),
        ("Calibré Afrique de l'Ouest",
         "Open source. Accessible pour une ONG à budget limité au Sahel."),
    ]
    for i, (title, body) in enumerate(diffs):
        y = 2.35 + i * 1.55
        box(s, 4.9, y, 0.06, 0.38, fill=AMB)
        tx(s, 5.1, y,      7.8, 0.42, title, sz=16, bold=True, col=WHT)
        tx(s, 5.1, y+0.45, 7.8, 0.6,  body,  sz=13, col=MUT)

    # Chiffre-clé discret bas droite
    tx(s, 5.1, 7.05, 7.8, 0.35,
       "1 évacuation évitée = 50× le coût annuel de l'abonnement",
       sz=11, col=MUT, it=True)

def s7_conclusion(p):
    s = slide(p)
    sn(s, 7)
    # Moitié gauche — capture STT/Alertes
    capture_ph(s, 0.4, 0.4, 7.5, 6.5,
               "Onglet Alertes ou TERROIR STT",
               "Montrer : scores par département, niveaux d'alerte, recommandations")
    # Moitié droite — message + démo
    vbar_x = 8.2
    box(s, vbar_x, 0.4, 0.06, 6.5, fill=GRN)
    # Tagline
    txm(s, 8.4, 0.7, 4.7, 2.0, [
        {'text': "Le territoire parle.", 'sz': 24, 'bold': True, 'col': WHT},
        {'text': "TERROIR écoute.",      'sz': 24, 'bold': True, 'col': GRN},
    ])
    # URL
    tx(s, 8.4, 2.9, 4.7, 0.45, URL, sz=13, col=GRN, bold=True)
    # QR code
    qr_ph(s, 8.4, 3.55, 2.0, 2.0)
    # Essayer maintenant
    tx(s, 10.55, 3.55, 2.5, 0.5,
       "Essayez\nmaintenant", sz=11, col=MUT, it=True)
    # Équipe
    box(s, 8.4, 6.5, 4.7, 0.04, fill=MUT)
    tx(s, 8.4, 6.62, 4.7, 0.45, TEAM + "  ·  " + "iSHEERO × DataCamp 2026",
       sz=11, col=MUT)

# ── MAIN ──────────────────────────────────────────────────────────────────────

def main():
    p = new_prs()
    s1_cover(p)
    s2_probleme(p)
    s3_outils(p)
    s4_carte(p)
    s5_comment(p)
    s6_impact(p)
    s7_conclusion(p)

    out = os.path.join(OUT, "TERROIR_pitch_deck_v5.pptx")
    p.save(out)
    print(f"PPTX sauvegardé : {out}")
    print(f"Slides générés  : {len(p.slides)}")

if __name__ == '__main__':
    main()
