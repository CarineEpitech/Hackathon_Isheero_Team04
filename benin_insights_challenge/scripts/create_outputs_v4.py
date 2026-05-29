#!/usr/bin/env python3
# Génère TERROIR_pitch_deck_v4.pptx, pitch_terroir_v4.md, qa_terroir_v4.md
# dans C:/Users/carin/OneDrive/Downloads/
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = r"C:\Users\carin\OneDrive\Downloads"

# ══════════════════════════════════════════════════════════════════
# PALETTE
# ══════════════════════════════════════════════════════════════════
BG   = RGBColor(0x0F, 0x17, 0x2A)
BOX  = RGBColor(0x1E, 0x29, 0x3B)
BOX2 = RGBColor(0x0A, 0x0F, 0x1A)
GRN  = RGBColor(0x05, 0x96, 0x69)
GRN2 = RGBColor(0x34, 0xD3, 0x99)
AMB  = RGBColor(0xF5, 0x9E, 0x0B)
RED  = RGBColor(0xDC, 0x26, 0x26)
WHT  = RGBColor(0xFF, 0xFF, 0xFF)
LGR  = RGBColor(0xCB, 0xD5, 0xE1)
MGR  = RGBColor(0x64, 0x74, 0x8B)
CYN  = RGBColor(0x22, 0xD3, 0xEE)

# ══════════════════════════════════════════════════════════════════
# HELPERS PPTX
# ══════════════════════════════════════════════════════════════════

def new_prs():
    p = Presentation()
    p.slide_width  = In(13.33)
    p.slide_height = In(7.5)
    return p

def add_slide(prs, bg=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = bg or BG
    return s

def rect(s, l, t, w, h, fill=None, lc=None):
    sh = s.shapes.add_shape(1, In(l), In(t), In(w), In(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if lc:
        sh.line.color.rgb = lc; sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh

def tx(s, l, t, w, h, text, sz=15, bold=False, col=None,
       al=PP_ALIGN.LEFT, it=False):
    col = col or LGR
    b  = s.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = b.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = al
    r  = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = it;   r.font.color.rgb = col
    return b

def txm(s, l, t, w, h, lines):
    """lines = list of dict {text,sz,bold,col,al,it} or empty string."""
    b  = s.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = b.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        if first: para = tf.paragraphs[0]; first = False
        else:     para = tf.add_paragraph()
        if not item:
            continue
        if isinstance(item, str):
            item = {"text": item}
        para.alignment = item.get("al", PP_ALIGN.LEFT)
        r = para.add_run()
        r.text          = item.get("text", "")
        r.font.size     = Pt(item.get("sz",   15))
        r.font.bold     = item.get("bold",  False)
        r.font.italic   = item.get("it",    False)
        r.font.color.rgb = item.get("col", LGR)
    return b

def top_bar(s, col=None):
    rect(s, 0, 0, 13.33, 0.07, fill=col or GRN)

def snum(s, n, total=11):
    tx(s, 12.5, 7.15, 0.7, 0.25, f"{n} / {total}", sz=9, col=MGR, al=PP_ALIGN.RIGHT)

def main_title(s, text, y=0.18, sz=27, col=None, al=PP_ALIGN.LEFT):
    tx(s, 0.6, y, 12.1, 1.05, text, sz=sz, bold=True, col=col or WHT, al=al)

def stat_box(s, l, t, w, h, value, label, vc=None, vsz=34, lsz=11):
    rect(s, l, t, w, h, fill=BOX)
    tx(s, l+0.1, t+0.1,    w-0.2, h*0.55,
       value, sz=vsz, bold=True, col=vc or AMB, al=PP_ALIGN.CENTER)
    tx(s, l+0.1, t+h*0.58, w-0.2, h*0.38,
       label, sz=lsz, col=LGR, al=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════

def s01_cover(prs):
    s = add_slide(prs)
    top_bar(s)
    # Ligne décorative
    rect(s, 0.6, 1.55, 1.5, 0.06, fill=GRN)
    # Titre principal
    tx(s, 0.6, 1.7, 12.1, 1.5, "TERROIR", sz=72, bold=True, col=WHT,
       al=PP_ALIGN.LEFT)
    # Sous-titre
    tx(s, 0.6, 3.3, 10.0, 0.7,
       "Veille territoriale et alerte précoce — Bénin, temps réel",
       sz=22, col=LGR, al=PP_ALIGN.LEFT)
    # Tagline
    tx(s, 0.6, 4.2, 9.0, 0.6,
       "Le territoire parle. TERROIR écoute.",
       sz=17, col=GRN2, it=True, al=PP_ALIGN.LEFT)
    # Séparateur bas
    rect(s, 0.6, 5.25, 12.1, 0.03, fill=MGR)
    # Info équipe
    tx(s, 0.6, 5.45, 12.1, 0.4,
       "BeninScope  ·  Team 04  ·  iSHEERO × DataCamp 2026",
       sz=12, col=MGR, al=PP_ALIGN.LEFT)
    tx(s, 0.6, 5.9, 12.1, 0.4,
       "Yaoitcha Rosine  ·  Gbaguidi Adewale  ·  Houndofi Jacques  ·  Agboton Carine",
       sz=11, col=MGR, al=PP_ALIGN.LEFT)
    # Pastilles décoratives (carte symbolique)
    rect(s, 11.5, 1.8, 0.18, 0.18, fill=GRN)
    rect(s, 11.9, 2.5, 0.14, 0.14, fill=AMB)
    rect(s, 11.3, 3.1, 0.12, 0.12, fill=RED)

def s02_probleme(prs):
    s = add_slide(prs)
    top_bar(s)
    snum(s, 2)
    main_title(s, "L'information était là. Personne ne la regardait.", sz=25)

    # Colonne gauche — 3 stat boxes
    stat_box(s,  0.6, 1.35, 3.8, 1.55, "4 %",   "couverture médiatique — nord Bénin",
             vc=AMB, vsz=42)
    stat_box(s,  0.6, 3.05, 3.8, 1.55, "−4,10", "ton médiatique nord (vs −1,09 sud)",
             vc=RED, vsz=42)
    stat_box(s,  0.6, 4.75, 3.8, 1.55, "0",     "système pour agréger ces signaux",
             vc=MGR, vsz=42)

    # Colonne droite — texte scénario
    rect(s, 4.7, 1.35, 8.0, 5.0, fill=BOX)
    txm(s, 4.9, 1.5, 7.7, 4.7, [
        {"text": "Le scénario concret",
         "sz": 14, "bold": True, "col": GRN},
        {"text": " ", "sz": 5, "col": LGR},
        {"text":
         "Un coordinateur ONG de santé communautaire à Kandi (Alibori) "
         "envoie ses agents motorisés chaque semaine vers des villages "
         "à 40–60 km. Avant chaque mission, il passe plusieurs appels "
         "pour « prendre la température ».",
         "sz": 13, "col": LGR},
        {"text": " ", "sz": 5, "col": LGR},
        {"text":
         "Cette veille informelle fonctionne — mais elle ne se transfère "
         "pas. Quand il s'absente, la connaissance part avec lui.",
         "sz": 13, "col": LGR},
        {"text": " ", "sz": 5, "col": LGR},
        {"text":
         "Les signaux sécuritaires existaient dans les médias régionaux "
         "(vols de bétail, tensions sur les axes, incidents de marché). "
         "Personne ne les lisait systématiquement.",
         "sz": 13, "col": LGR},
        {"text": " ", "sz": 7, "col": LGR},
        {"text": "Ce n'est pas un manque de compétence.",
         "sz": 14, "bold": True, "col": WHT},
        {"text": "C'est un manque de système.",
         "sz": 14, "bold": True, "col": AMB},
    ])

def s03_outils(prs):
    s = add_slide(prs)
    top_bar(s)
    snum(s, 3)
    main_title(s, "Ce qui existe ne répond pas à ce problème.", sz=26)

    # En-têtes colonnes
    rect(s, 0.6,  1.4,  5.8, 0.45, fill=BOX)
    rect(s, 6.85, 1.4,  6.0, 0.45, fill=BOX)
    tx(s, 0.7,  1.48, 5.6, 0.35, "Ce qui existe",    sz=13, bold=True, col=MGR)
    tx(s, 6.95, 1.48, 5.8, 0.35, "Ce qui manque",    sz=13, bold=True, col=GRN)

    rows = [
        ("ACLED — enregistre les incidents violents",
         "Signal avant l'incident — pas après"),
        ("Rapports institutionnels — mensuels",
         "Mise à jour continue — pas hebdomadaire"),
        ("Veille manuelle — impossible à tenir",
         "Agrégation automatique multi-sources"),
        ("Dataminr / ICEWS — millions d'euros, usage gouvernemental occidental",
         "Accessible pour une ONG sahélienne à budget limité"),
    ]
    for i, (left, right) in enumerate(rows):
        y = 1.95 + i * 1.28
        rect(s, 0.6,  y, 5.8, 1.1, fill=BOX)
        rect(s, 6.85, y, 6.0, 1.1, fill=BOX, lc=GRN)
        tx(s, 0.75,  y+0.12, 5.5, 0.85, left,  sz=13, col=LGR)
        tx(s, 7.0,   y+0.12, 5.7, 0.85, right, sz=13, col=GRN2)
    # Flèche centrale symbolique
    for i in range(4):
        y = 1.95 + i * 1.28 + 0.42
        tx(s, 6.45, y, 0.4, 0.3, "→", sz=16, bold=True, col=AMB, al=PP_ALIGN.CENTER)

def s04_phase1(prs):
    s = add_slide(prs)
    top_bar(s, col=AMB)
    snum(s, 4)
    main_title(s, "On a analysé 23 859 événements. Voici ce qu'on a trouvé.", sz=24)

    # 3 grandes stat boxes horizontales
    stat_box(s, 0.6,  1.4, 3.9, 2.2, "23 859", "événements GDELT\nBénin 2025", vc=CYN, vsz=38)
    stat_box(s, 4.7,  1.4, 3.9, 2.2, "4 × plus\nnégatif",
             "ton nord (−4,10) vs sud (−1,09)", vc=RED, vsz=28)
    stat_box(s, 8.8,  1.4, 3.9, 2.2, "25 %",
             "d'événements conflictuels\nstructurent toute la perception", vc=AMB, vsz=38)

    # Bloc insight + transition
    rect(s, 0.6, 3.8, 12.1, 1.4, fill=BOX)
    txm(s, 0.8, 3.95, 11.8, 1.2, [
        {"text": "Ce que ça révèle sur le nord :", "sz": 13, "bold": True, "col": GRN},
        {"text":
         "Quand le nord est couvert par les médias internationaux, c'est presque toujours "
         "pour un incident sécuritaire. Vol de bétail, tension sur un axe routier, conflit "
         "communautaire. Le signal existe — il n'est pas agrégé.",
         "sz": 13, "col": LGR},
    ])

    # Ligne de transition
    rect(s, 0.6, 5.35, 12.1, 0.06, fill=AMB)
    txm(s, 0.6, 5.5, 12.1, 0.9, [
        {"text":
         "La Phase 1 a confirmé que le signal existe.  "
         "La Phase 2 crée le système pour l'utiliser.",
         "sz": 17, "bold": True, "col": WHT, "al": PP_ALIGN.CENTER},
    ])

def s05_terroir(prs):
    s = add_slide(prs)
    top_bar(s)
    snum(s, 5)
    main_title(s, "Une carte live. Un score. Une alerte.", sz=28)

    # 3 colonnes + flèches
    cols = [
        (0.6,  "GDELT",         GRN,  ["Médias mondiaux", "Mise à jour toutes", "les 15 minutes"]),
        (4.85, "Signalements",  CYN,  ["Citoyens + ONG", "GPS géolocalisé", "Temps réel"]),
        (9.1,  "TERROIR",       AMB,  ["Score / Carte / Alerte", "Par département", "En continu"]),
    ]
    for l, title_txt, col, items in cols:
        rect(s, l, 1.35, 3.6, 3.8, fill=BOX)
        rect(s, l, 1.35, 3.6, 0.45, fill=col)
        tx(s, l+0.12, 1.43, 3.36, 0.35, title_txt, sz=14, bold=True, col=BG, al=PP_ALIGN.CENTER)
        y = 2.05
        for item in items:
            tx(s, l+0.12, y, 3.36, 0.5, item, sz=14, col=LGR, al=PP_ALIGN.CENTER)
            y += 0.55
    # Flèches +  →
    tx(s, 4.3,  2.85, 0.6, 0.5, "+", sz=28, bold=True, col=GRN2, al=PP_ALIGN.CENTER)
    tx(s, 8.55, 2.85, 0.6, 0.5, "→", sz=28, bold=True, col=AMB,  al=PP_ALIGN.CENTER)

    # Ligne basse
    rect(s, 0.6, 5.35, 12.1, 0.06, fill=GRN)
    tx(s, 0.6, 5.55, 12.1, 0.65,
       "Tout utilisateur — citoyen, ONG, mairie, journaliste — voit la même carte "
       "en temps réel, sans inscription obligatoire.",
       sz=15, col=LGR, al=PP_ALIGN.CENTER)
    # Sous-accroche
    txm(s, 0.6, 6.3, 12.1, 0.7, [
        {"text":
         "Un signalement citoyen à Kandi est visible en quelques secondes "
         "par tous les coordinateurs connectés au nord du Bénin.",
         "sz": 13, "it": True, "col": MGR, "al": PP_ALIGN.CENTER},
    ])

def s06_stt(prs):
    s = add_slide(prs)
    top_bar(s)
    snum(s, 6)
    main_title(s, "Le Score de Tension Territorial", sz=28)

    # Formule simplifiée
    rect(s, 0.6, 1.35, 12.1, 1.05, fill=BOX)
    txm(s, 0.8, 1.45, 11.8, 0.9, [
        {"text": "STT  =  ",
         "sz": 16, "bold": True, "col": WHT},
        {"text":
         "40 % codes conflictuels  +  35 % ton négatif  +  15 % volume  +  10 % diversité sources",
         "sz": 14, "col": LGR},
        {"text": "  Fenêtre 14 jours  ·  Baseline 90 jours  ·  Par département  ·  Calculé quotidiennement",
         "sz": 11, "it": True, "col": MGR},
    ])

    # 2 niveaux d'alerte
    rect(s, 0.6, 2.55, 5.85, 1.0, fill=BOX, lc=AMB)
    rect(s, 6.65, 2.55, 6.05, 1.0, fill=BOX, lc=RED)
    txm(s, 0.75, 2.65, 5.6, 0.8, [
        {"text": "Niveau 1 — Précaution  (STT ≥ 2,0)", "sz": 13, "bold": True, "col": AMB},
        {"text": "Vérification terrain recommandée dans 48h", "sz": 12, "col": LGR},
    ])
    txm(s, 6.8, 2.65, 5.8, 0.8, [
        {"text": "Niveau 2 — Alerte  (STT ≥ 3,0)", "sz": 13, "bold": True, "col": RED},
        {"text": "Notification OCHA + évaluation repositionnement", "sz": 12, "col": LGR},
    ])

    # Boîte d'alerte exemple
    rect(s, 0.6, 3.7, 12.1, 3.05, fill=BOX2)
    rect(s, 0.6, 3.7, 0.08, 3.05, fill=RED)
    txm(s, 0.82, 3.82, 11.7, 2.8, [
        {"text": "EXEMPLE D'ALERTE — Alibori",
         "sz": 12, "bold": True, "col": RED},
        {"text": "Score : 7,2 / 10   [Niveau Précaution — seuil 6,0]",
         "sz": 13, "bold": True, "col": AMB},
        {"text": " ", "sz": 5, "col": LGR},
        {"text": "Signaux détectés :",
         "sz": 12, "bold": True, "col": LGR},
        {"text": "— Incidents sécuritaires CAMEO (vols, tensions marchés) : +78 % sur 14 jours",
         "sz": 12, "col": LGR},
        {"text": "— Ton moyen : −4,1  (vs −1,8 médiane nationale)",
         "sz": 12, "col": LGR},
        {"text": "— 5 médias convergent sur la zone (RFI, BBC Afrique, VOA Afrique, Sahel Intel., Lefaso)",
         "sz": 12, "col": LGR},
        {"text": " ", "sz": 5, "col": LGR},
        {"text": "Précédent le plus proche : mars 2024  (score 6,8 → restriction d'accès 7 jours plus tard)",
         "sz": 11, "it": True, "col": MGR},
        {"text": "Action recommandée : vérification terrain dans les 48h — si confirmé → notifier OCHA",
         "sz": 12, "bold": True, "col": GRN},
    ])

def s07_demo(prs):
    s = add_slide(prs)
    top_bar(s)
    snum(s, 7)
    main_title(s, "Ce que vit le coordinateur à Kandi (Alibori)", sz=25)

    steps = [
        (AMB, "J0  matin",
         "Alerte TERROIR reçue.",
         "Score 7,2/10 — Alibori.\nSignaux : vols bétail +78 %,\nton −4,1."),
        (CYN, "J0  après-midi",
         "Contact réseau local.",
         "Correspondant confirme :\ntension sur l'axe nord.\nRecommandation vérifiée."),
        (GRN, "J1",
         "Décision documentée.",
         "Mission reportée de 48h.\nRéseau terrain prévenu.\nDécision tracée."),
        (GRN2, "Résultat",
         "Zéro incident.",
         "Zéro urgence.\nDécision prise avec\nun signal, pas du bruit."),
    ]
    for i, (col, step, action, detail) in enumerate(steps):
        l = 0.6 + i * 3.2
        # Boîte
        rect(s, l, 1.4, 2.95, 4.8, fill=BOX)
        # Bande couleur haut
        rect(s, l, 1.4, 2.95, 0.4, fill=col)
        # Étape
        tx(s, l+0.08, 1.45, 2.79, 0.32, step, sz=12, bold=True, col=BG, al=PP_ALIGN.CENTER)
        # Action
        tx(s, l+0.1, 1.95, 2.75, 0.5, action, sz=14, bold=True, col=WHT)
        # Détail
        tx(s, l+0.1, 2.55, 2.75, 2.5, detail, sz=13, col=LGR)
    # Flèches
    for i in range(3):
        tx(s, 3.5 + i*3.2, 3.5, 0.35, 0.5, "→", sz=18, bold=True, col=MGR, al=PP_ALIGN.CENTER)
    # Encadré bas
    rect(s, 0.6, 6.35, 12.1, 0.85, fill=BOX)
    tx(s, 0.75, 6.48, 11.8, 0.6,
       "La décision reste humaine. TERROIR fournit le signal. Le coordinateur évalue, "
       "décide, documente. Le système apprend du retour terrain.",
       sz=13, col=LGR)

def s08_valeur(prs):
    s = add_slide(prs)
    top_bar(s, col=AMB)
    snum(s, 8)
    main_title(s, "Le calcul est simple.", sz=30)

    # 3 blocs empilés
    blocs = [
        (RED,  "50 000 – 200 000 €",  "coût d'une évacuation précipitée d'équipes terrain", 38),
        (GRN,  "quelques milliers €/an",  "abonnement annuel à TERROIR (TerritoryPulse)", 30),
        (AMB,  "1 évacuation évitée",  "= 50 × le coût de l'abonnement", 30),
    ]
    for i, (col, value, label, vsz) in enumerate(blocs):
        y = 1.4 + i * 1.6
        rect(s, 0.6, y, 0.12, 1.35, fill=col)
        rect(s, 0.85, y, 11.85, 1.35, fill=BOX)
        tx(s, 1.0,  y+0.08, 7.0, 0.7, value, sz=vsz, bold=True, col=col)
        tx(s, 1.0,  y+0.78, 10.5, 0.45, label, sz=14, col=LGR)

    # Ligne marché
    rect(s, 0.6, 6.15, 12.1, 0.65, fill=BOX)
    tx(s, 0.75, 6.25, 11.8, 0.5,
       "20 ONGs actives au Sahel et sur la côte ouest-africaine  —  "
       "marché immédiat : 60 000 – 160 000 €/an",
       sz=13, col=MGR)

def s09_diff(prs):
    s = add_slide(prs)
    top_bar(s)
    snum(s, 9)
    main_title(s, "TERROIR surveille. Vous décidez.", sz=28)

    # Colonne gauche : 4 arguments
    args = [
        (GRN,  "Actif, pas passif",
         "TERROIR envoie l'alerte. "
         "Vous n'avez pas besoin d'ouvrir l'application."),
        (AMB,  "Avant, pas après",
         "ACLED documente ce qui s'est passé. "
         "TERROIR signale ce qui commence à se passer."),
        (CYN,  "Calibré pour l'Afrique de l'Ouest",
         "ICEWS et Dataminr sont construits pour les gouvernements occidentaux. "
         "TERROIR est fait pour Kandi, Parakou, Natitingou."),
        (GRN2, "Il apprend",
         "Chaque alerte génère un retour terrain. "
         "Les faux positifs recalibrent les seuils. "
         "Le système s'améliore mois après mois."),
    ]
    for i, (col, title_txt, body) in enumerate(args):
        y = 1.4 + i * 1.38
        rect(s, 0.6, y, 0.08, 1.2, fill=col)
        rect(s, 0.82, y, 7.3, 1.2, fill=BOX)
        tx(s, 0.97, y+0.08, 7.0, 0.38, title_txt, sz=14, bold=True, col=col)
        tx(s, 0.97, y+0.5, 7.0, 0.6, body, sz=12, col=LGR)

    # Tableau comparatif droite
    headers = ["", "TERROIR", "ACLED", "Dashboard\npassif"]
    rows_data = [
        ("Fréquence",   "15 min",       "Hebdo/retard",  "Manuel"),
        ("Direction",   "Avant",        "Après",         "Aucune"),
        ("Alerte auto", "Oui",          "Non",           "Non"),
        ("Accessibilité","ONG sahélienne","Chercheurs",   "Quiconque"),
    ]
    col_w = [2.0, 1.55, 1.55, 1.55]
    col_x = [8.3, 10.3, 11.85, 13.4]
    # Oops - let me recalculate to fit 13.33 width
    # Right column starts at x=8.6, total width = 4.6
    # 4 columns in 4.6 inches
    TBX = 8.62
    TBW = 4.6
    col_ws = [1.35, 1.1, 1.1, 1.05]
    col_xs = [TBX, TBX+1.35, TBX+2.45, TBX+3.55]

    # Header
    for j, (hd, cw, cx) in enumerate(zip(headers, col_ws, col_xs)):
        fill = GRN if j == 1 else BOX
        col_hd = BG if j == 1 else MGR
        rect(s, cx, 1.4, cw-0.03, 0.45, fill=fill)
        tx(s, cx+0.04, 1.47, cw-0.1, 0.35,
           hd, sz=11, bold=True, col=col_hd if j==1 else WHT,
           al=PP_ALIGN.CENTER)

    for i, (lbl, v1, v2, v3) in enumerate(rows_data):
        y = 1.95 + i * 1.1
        vals = [lbl, v1, v2, v3]
        for j, (val, cw, cx) in enumerate(zip(vals, col_ws, col_xs)):
            fill = BOX if j != 1 else RGBColor(0x0D, 0x28, 0x1C)
            col_txt = LGR if j != 1 else GRN2
            rect(s, cx, y, cw-0.03, 0.95, fill=fill)
            tx(s, cx+0.05, y+0.1, cw-0.12, 0.75,
               val, sz=11, col=col_txt, al=PP_ALIGN.CENTER)

def s10_limits(prs):
    s = add_slide(prs)
    top_bar(s, col=AMB)
    snum(s, 10)
    main_title(s, "Ce qu'on sait. Ce qu'on construit.", sz=27)

    # En-têtes des 3 colonnes
    headers = [
        (RED,   "Limite connue"),
        (GRN,   "Ce qu'on a fait"),
        (CYN,   "Ce qui vient"),
    ]
    col_x = [0.6, 4.65, 8.7]
    col_w = 3.85
    for (col, hd), x in zip(headers, col_x):
        rect(s, x, 1.4, col_w, 0.42, fill=col)
        tx(s, x+0.12, 1.47, col_w-0.2, 0.3,
           hd, sz=12, bold=True, col=BG, al=PP_ALIGN.CENTER)

    # 3 lignes de contenu
    rows = [
        (
            "91 % des événements GDELT géolocalisés au niveau pays — pas de précision communale",
            "Couche signalement citoyen GPS : correction directe du biais géographique",
            "Intégration ACLED pour résolution ADM2 (commune) dans les 6 mois",
        ),
        (
            "GDELT mesure les médias internationaux — pas la réalité terrain",
            "Double flux : signal médiatique + signalements terrain validés avec statut explicite",
            "Extension réseau correspondants locaux ONG et journalistes nord-Bénin",
        ),
        (
            "Calibration STT jeune — moins d'un an de données de référence",
            "Boucle de feedback intégrée dès le lancement : chaque alerte génère un retour",
            "6 mois de données → affinage automatique des seuils par département",
        ),
    ]
    for i, (lim, fait, vient) in enumerate(rows):
        y = 1.95 + i * 1.6
        for j, (text, x) in enumerate(zip([lim, fait, vient], col_x)):
            border_col = [RED, GRN, CYN][j]
            rect(s, x, y, col_w, 1.45, fill=BOX, lc=border_col)
            tx(s, x+0.14, y+0.12, col_w-0.26, 1.2, text, sz=12, col=LGR)

    # Ligne de principe bas
    rect(s, 0.6, 6.7, 12.1, 0.55, fill=BOX)
    tx(s, 0.75, 6.79, 11.8, 0.38,
       "Un système d'alerte qui documente ses limites est plus crédible "
       "qu'un système qui les cache. C'est la même logique qu'ACLED.",
       sz=12, it=True, col=MGR)


def s11_conclusion(prs):
    s = add_slide(prs)
    top_bar(s)
    # Tagline centrale
    tx(s, 0.6, 0.6, 12.1, 1.1,
       "Le territoire parle. TERROIR écoute.",
       sz=30, bold=True, col=WHT, al=PP_ALIGN.LEFT)
    tx(s, 0.6, 1.7, 9.0, 0.55,
       "Et distribue l'information à ceux qui en ont besoin, "
       "avant qu'il soit trop tard pour agir.",
       sz=16, col=LGR, it=True)

    # 4 blocs "Ce qu'on a construit"
    items = [
        (GRN,  "Plateforme opérationnelle",
         "Streamlit · GDELT live · SQLite\nCarte · Alertes · Signalement"),
        (AMB,  "Score de Tension Territorial",
         "Par département · Quotidien\nZ-scores · Baseline 90 jours"),
        (CYN,  "Signalement citoyen",
         "GPS · Temps réel · Modération\nBoucle de calibration terrain"),
        (GRN2, "Chaîne complète",
         "Détection → Score → Alerte\n→ Recommandation → Décision"),
    ]
    for i, (col, ttl_txt, body) in enumerate(items):
        l = 0.6 + i * 3.2
        rect(s, l, 2.45, 2.95, 2.8, fill=BOX)
        rect(s, l, 2.45, 2.95, 0.38, fill=col)
        tx(s, l+0.1, 2.5, 2.75, 0.3, ttl_txt, sz=11, bold=True, col=BG, al=PP_ALIGN.CENTER)
        tx(s, l+0.1, 2.95, 2.75, 2.15, body, sz=12, col=LGR, al=PP_ALIGN.CENTER)

    # Prochaine étape
    rect(s, 0.6, 5.45, 12.1, 0.6, fill=BOX)
    txm(s, 0.75, 5.54, 11.8, 0.45, [
        {"text": "Prochaine étape : ",
         "sz": 14, "bold": True, "col": AMB},
        {"text": "Burkina Faso  ·  Niger  ·  Partenariat clusters sécurité OCHA",
         "sz": 14, "col": LGR},
    ])

    # GitHub
    tx(s, 0.6, 6.25, 8.0, 0.4,
       "github.com/CarineEpitech/Hackathon_Isheero_Team04",
       sz=11, col=MGR)
    snum(s, 11)

# ══════════════════════════════════════════════════════════════════
# GÉNÉRATION PPTX
# ══════════════════════════════════════════════════════════════════

def build_pptx():
    prs = new_prs()
    s01_cover(prs)
    s02_probleme(prs)
    s03_outils(prs)
    s04_phase1(prs)
    s05_terroir(prs)
    s06_stt(prs)
    s07_demo(prs)
    s08_valeur(prs)
    s09_diff(prs)
    s10_limits(prs)
    s11_conclusion(prs)
    out = os.path.join(OUT, "TERROIR_pitch_deck_v4.pptx")
    prs.save(out)
    print(f"PPTX : {out}")

# ══════════════════════════════════════════════════════════════════
# PITCH MARKDOWN
# ══════════════════════════════════════════════════════════════════

PITCH_MD = """\
# Scripts de pitch TERROIR — v4

---

## Pitch A — Simple et humain (3 minutes)

---

Le nord du Bénin représente quatre pour cent de la couverture médiatique internationale du pays. Mais quand il est couvert, le ton est quatre fois plus négatif que le reste. Vol de bétail, tensions sur les axes routiers, incidents de marché. Les signaux sont là. Personne ne les lit systématiquement.

Un coordinateur ONG de santé communautaire à Kandi envoie ses agents terrain chaque semaine. Avant chaque mission, il passe plusieurs appels pour prendre la température. Ce réseau informel fonctionne — mais il ne se transfère pas. Quand il part en congé, la connaissance part avec lui.

Ce n'est pas un problème d'information. C'est un manque de système.

---

En Phase 1, on a analysé 23 859 événements GDELT couvrant toute l'année 2025. Ce qu'on a confirmé : les signaux sécuritaires précèdent les incidents. Ils sont mesurables. Ils sont dans les données. Ils ne sont simplement pas agrégés pour être utiles.

TERROIR, c'est notre réponse.

Deux sources d'information : les médias internationaux via GDELT, mis à jour toutes les quinze minutes. Et les signalements directs des citoyens et des agents terrain, géolocalisés en temps réel sur une carte partagée.

On combine ces deux flux. On calcule un Score de Tension pour chaque département du Bénin — chaque jour, en continu. Quand ce score dépasse un seuil dans l'Alibori ou le Borgou, une alerte part automatiquement. Zone, score, signaux, précédent, recommandation. En quelques lignes.

Le coordinateur à Kandi n'a plus besoin de passer dix appels le matin. TERROIR lui donne le signal. La décision reste la sienne.

---

La valeur est directe. Une mission annulée à tort coûte du temps. Un incident évitable coûte bien plus. Un abonnement à TERROIR coûte quelques milliers d'euros par an. La première décision mieux éclairée rembourse l'abonnement.

Ce n'est pas un tableau de bord qu'on consulte quand on y pense. C'est un système qui surveille en permanence et qui alerte quand c'est nécessaire.

ACLED documente ce qui s'est passé. Dataminr sert les gouvernements occidentaux. TERROIR est construit pour les acteurs terrain en Afrique de l'Ouest, sur des données ouvertes, pour le prix d'un abonnement d'ONG.

---

Le territoire parle. Les citoyens aussi.
TERROIR écoute — et distribue l'information à ceux qui en ont besoin, avant qu'il soit trop tard pour agir.

---

## Pitch B — Jury technique (3 minutes)

---

Notre point de départ : 23 859 événements GDELT couvrant le Bénin sur 2025. L'analyse Phase 1 a révélé un paradoxe territorial : le nord représente quatre pour cent de la couverture internationale — mais son ton médiatique est à moins 4,10, contre moins 1,09 pour le sud. Quand le nord est couvert, c'est presque exclusivement pour des incidents de sécurité interne — vols de bétail, tensions routières, conflits de marché. Des codes CAMEO négatifs sur un volume faible mais structurellement significatif.

Ces signaux ont une valeur opérationnelle directe. À condition de les automatiser.

---

TERROIR transforme ce corpus analytique en système d'alerte quasi temps réel. Architecture en deux flux.

Flux un : GDELT filtré sur les codes CAMEO sécuritaires, latence quinze minutes, géolocalisé au niveau ADM1 pour les départements prioritaires — Alibori, Borgou, Atacora.

Flux deux : signalement communautaire. Formulaire géolocalisé GPS, enregistrement SQLite instantané, affichage carte live avec statut de vérification explicite — validé ou non vérifié.

---

Les deux flux alimentent le Score de Tension Territorial. Combinaison pondérée de quatre z-scores : part d'événements CAMEO négatifs — quarante pour cent —, ton moyen inversé — trente-cinq —, volume normalisé — quinze —, entropie de Shannon sur les sources — dix. Fenêtre courante de quatorze jours, baseline de quatre-vingt-dix jours. Deux niveaux d'alerte : Précaution à 2,0, Alerte à 3,0.

La boucle de calibration est intégrée. Chaque alerte génère un retour terrain — confirmé, non confirmé, partiel — qui recalibre les seuils sur la zone. Le taux de faux positifs se réduit mécaniquement sur six mois.

---

Deux limites assumées publiquement : 91,2 % des événements GDELT ont une géolocalisation générique niveau pays. C'est le biais structurel documenté de GDELT. La couche de signalement citoyen le corrige directement — c'est sa fonction première, pas un ajout de confort.

---

La différenciation est contextuelle. ACLED documente a posteriori. ICEWS et Dataminr sont calibrés pour des dynamiques médiatiques occidentales, inaccessibles pour une ONG sahélienne. TERROIR est construit sur des données ouvertes, code auditable, pour un usage spécifiquement ouest-africain.

Marché immédiat : coordinateurs sécurité des ONGs internationales au nord Bénin et au Sahel. Modèle : abonnement annuel. Le ROI se comprend à la première décision mieux éclairée.

---

Le territoire parle. TERROIR rend ce signal lisible, actionnable, et distribué en temps réel.
"""

# ══════════════════════════════════════════════════════════════════
# Q&A MARKDOWN
# ══════════════════════════════════════════════════════════════════

QA_MD = """\
# Réponses aux questions de soumission — TERROIR v4

---

## 01 / Problème — Le problème adressé

### Version A — Synthétique

Le nord du Bénin concentre 4 % de la couverture médiatique internationale du pays — mais son ton est quatre fois plus négatif que le reste (−4,10 vs −1,09). Quand cette zone est couverte, c'est presque exclusivement pour des incidents de sécurité interne : vols de bétail, tensions sur les axes routiers commerciaux, conflits communautaires liés au foncier, incidents de marché frontalier.

Ces signaux existent. Ils ne sont pas agrégés. Les acteurs qui opèrent dans ces zones — coordinateurs d'ONG, agents préfectoraux, agents de terrain — prennent leurs décisions sur la base du bouche-à-oreille et de l'expérience personnelle, faute d'un système capable de consolider ces signaux au moment où ils comptent.

---

### Version B — Standard

L'analyse de 23 859 événements GDELT couvrant le Bénin sur l'année 2025 a révélé un paradoxe territorial : le nord du pays est à la fois le moins couvert par les médias internationaux et le plus négatif sur le plan sécuritaire. Alibori, Borgou, Atacora — ces départements représentent 4 % de la couverture, mais leur ton médiatique est quatre fois plus négatif que le sud, dominé par des incidents de sécurité interne quotidienne.

La sécurité interne dans le nord du Bénin ne se réduit pas aux conflits armés. Elle recouvre un spectre plus large : vols de bétail récurrents dans les zones rurales de l'Alibori, insécurité sur les axes routiers commerciaux (Parakou–Malanville, Natitingou–Tanguiéta), tensions communautaires entre agriculteurs et éleveurs à recrudescence saisonnière, agressions sur les agents de terrain.

Ces incidents suivent des cycles mesurables — saisonniers, économiques, liés aux périodes de marché. Mais aucun acteur ne dispose d'un système pour les agréger, les contextualiser, et les distribuer au bon moment à la bonne personne. Les décisions de terrain se prennent sur la base du bouche-à-oreille, pas d'un signal consolidé.

Le problème n'est pas un manque d'information. C'est un manque de système pour la rendre utile.

---

### Version C — Complète

L'analyse de Phase 1 a confirmé statistiquement ce que les acteurs terrain savent intuitivement : le nord du Bénin est une zone à forte vulnérabilité sécuritaire, sous-couverte par les médias internationaux, et dont les signaux précurseurs ne sont pas exploités.

Quatre pour cent de la couverture médiatique du pays. Ton à −4,10 — soit quatre fois plus négatif que le sud. Et quand on décompose les codes CAMEO, la dominante est claire : tensions, coercition, incidents de basse intensité. Pas de grands conflits armés — mais une insécurité interne quotidienne qui pèse directement sur les décisions des acteurs opérationnels.

Vols de bétail récurrents dans l'Alibori. Insécurité sur l'axe Parakou–Malanville. Tensions communautaires entre agriculteurs et éleveurs à recrudescence saisonnière. Agressions sur les agents de terrain en moto vers les villages enclavés. Pressions sur les marchés hebdomadaires frontaliers.

Ces incidents ont une caractéristique commune : ils se produisent loin des centres, ils ne remontent que rarement dans les médias nationaux ou internationaux, et ils suivent des cycles que personne n'a encore documentés systématiquement.

Un coordinateur ONG à Kandi envoie ses agents terrain chaque semaine. Avant chaque mission, il passe plusieurs appels pour évaluer la situation. Cette veille informelle fonctionne — mais elle ne se transfère pas, ne laisse aucune trace, et disparaît quand il part en congé. Ce n'est pas un manque de compétence. C'est un manque de système.

Les outils existants ne comblent pas ce vide : ACLED enregistre après l'incident, les rapports institutionnels sont hebdomadaires ou mensuels, les systèmes commerciaux coûtent des millions et sont calibrés pour des dynamiques médiatiques occidentales. Aucun n'est construit pour le nord du Bénin.

---

## 02 / Solution — L'approche proposée

### Version A — Synthétique

TERROIR est une plateforme de veille territoriale quasi temps réel pour le Bénin. Elle agrège les événements GDELT (mis à jour toutes les 15 minutes) et les signalements citoyens géolocalisés, calcule un Score de Tension Territorial par département, et envoie automatiquement une alerte structurée quand un seuil est franchi — avec zone, score, signaux déclencheurs et recommandation d'action.

Elle est calibrée sur les dynamiques sécuritaires de l'Afrique de l'Ouest — pas sur les modèles des systèmes occidentaux.

---

### Version B — Standard

TERROIR agrège deux flux complémentaires : les événements GDELT filtrés sur les codes sécuritaires du Bénin (latence 15 minutes), et les signalements citoyens géolocalisés par GPS, visibles en temps réel sur une carte partagée.

Ces deux flux alimentent un Score de Tension Territorial calculé quotidiennement pour chaque département. Ce score compare la fenêtre courante de 14 jours à une baseline de 90 jours sur quatre dimensions : part d'événements conflictuels (40 %), ton de la couverture (35 %), volume normalisé (15 %), diversité des sources (10 %).

Quand le score dépasse un seuil, une alerte est envoyée automatiquement. Zone, score, signaux déclencheurs, précédent historique, recommandation opérationnelle — le message est court et actionnable.

La plateforme est accessible sans inscription pour la consultation. Le code est auditable. Les données sont ouvertes. Elle est construite pour les acteurs terrain en Afrique de l'Ouest.

---

### Version C — Développée

TERROIR (Tableau d'Évaluation et de Remontée Rapide des Observations et Incidents sur les Risques) est une plateforme citoyenne de veille de sécurité territoriale en quasi temps réel pour le Bénin.

**Architecture à deux flux :**

Flux 1 — GDELT automatisé : événements ingérés toutes les quinze minutes, filtrés sur les codes CAMEO sécuritaires et conflictuels, géolocalisés au niveau ADM1. Latence inférieure à trente minutes entre la publication d'un article et son apparition sur la carte.

Flux 2 — Signalement communautaire : formulaire géolocalisé GPS (type d'incident, description libre, position auto-détectée). Enregistrement SQLite instantané. Affichage sur carte avec statut de vérification explicite. Recoupement automatique avec GDELT pour validation.

**Score de Tension Territorial (STT) :**
Combinaison pondérée de quatre z-scores sur fenêtre glissante de 14 jours / baseline 90 jours :
- Part d'événements CAMEO négatifs : 40 %
- Ton moyen inversé : 35 %
- Volume normalisé : 15 %
- Entropie de Shannon sur les sources : 10 %

Deux niveaux d'alerte : Précaution à STT ≥ 2,0 / Alerte à STT ≥ 3,0.

**Boucle de calibration :** chaque alerte génère un retour terrain (confirmé / non confirmé / partiel) qui recalibre les seuils sur la zone. Le taux de faux positifs se réduit mécaniquement sur six mois.

**Stack technique :** Python · Streamlit · GDELT BigQuery · SQLite · Folium · Plotly.

**Limites assumées :** 91,2 % des événements GDELT ont une géolocalisation générique niveau pays. La couche de signalement citoyen corrige directement cette limite — c'est sa raison d'être première. Le jury est informé de cette limite dans l'interface même.

**Différenciation :** TERROIR ne prédit pas — il signale des anomalies mesurables et actionnables. Il travaille en amont d'ACLED. Il est calibré pour les médias ouest-africains. Il est accessible pour une ONG à budget limité au nord Bénin.

---

## 03 / L'équipe

**Team 04 — iSHEERO × DataCamp 2026**

- **Agboton Carine** — Architecture de données, pipeline GDELT, coordination technique
- **Gbaguidi Adewale** — Développement backend, plateforme Streamlit, déploiement
- **Houndofi Jacques** — Analyse et modélisation, Score de Tension Territorial, méthodes statistiques
- **Yaoitcha Rosine** — Visualisation, graphiques interactifs, structuration du pitch

Quatre étudiants béninois qui ont choisi d'aller au-delà de l'analyse pour construire quelque chose d'utile au Bénin.
"""

# ══════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════

def write_md(filename, content):
    path = os.path.join(OUT, filename)
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"MD   : {path}")

def main():
    os.makedirs(OUT, exist_ok=True)
    build_pptx()
    write_md("pitch_terroir_v4.md", PITCH_MD)
    write_md("qa_terroir_v4.md",    QA_MD)
    print("\nTerminé. 3 fichiers générés dans :", OUT)

if __name__ == "__main__":
    main()
